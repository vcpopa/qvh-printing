# pylint: disable=line-too-long
"""
Module: Presentation Utilities

This module provides functions for working with PowerPoint presentations (PPTX files).
It includes utilities for copying slides between presentations and merging slides from
multiple presentations into a single presentation.

Functions:
- copy_slide(slide, new_presentation):
    Copy a slide from one presentation to another while preserving formatting.

- merge_presentations(directory_path, output_filename):
    Merge slides from multiple presentations in a directory into one.

Dependencies:
- os: Provides functions for interacting with the operating system, used for file operations.
- copy: Provides functions for shallow and deep copy operations.
- pptx.Presentation: Represents a PowerPoint presentation and provides methods to manipulate slides.
- pptx.enum.shapes.MSO_SHAPE_TYPE: Enumerates the types of shapes in PowerPoint slides.

Example usage:
    from ppt_utils import copy_slide, merge_presentations

    # Copy a slide from one presentation to another
    source_presentation = Presentation('source.pptx')
    new_presentation = Presentation()
    slide_to_copy = source_presentation.slides[0]
    copy_slide(slide_to_copy, new_presentation)

    # Merge slides from a directory into a single presentation
    merge_presentations('slides_directory', 'merged_output.pptx')
"""
import os
import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

PPT_WIDTH=13.33
PPT_HEIGHT=7.5

def copy_slide(slide, new_presentation):
    """
    Copy a slide from one presentation to another while preserving formatting.

    Parameters:
    - slide (Slide): Slide object to copy.
    - new_presentation (Presentation): Presentation object to copy the slide into.

    Returns:
    - None
    """
    img_dict = {}
    layout = new_presentation.slide_masters[0].slide_layouts[
        6
    ]  # Change the index according to your desired layout
    new_slide = new_presentation.slides.add_slide(layout)

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Generate a unique name for the picture file
            img_filename = f"image_{len(img_dict)}.jpg"
            img_dict[img_filename] = [
                shape.left,
                shape.top,
                shape.width,
                shape.height,
            ]
            # Save the image to a file
            with open(img_filename, "wb") as f:
                f.write(shape.image.blob)
        else:
            # For other shapes, copy the shape element
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Add pictures from the image dictionary to the new slide
    for img_path, position in img_dict.items():
        new_slide.shapes.add_picture(
            img_path, position[0], position[1], position[2], position[3]
        )

    # Remove the temporary image files
    for img_path in img_dict:
        os.remove(img_path)


def merge_presentations(directory_path, output_filename):
    """
    Merge slides from multiple presentations in a directory into one.

    Parameters:
    - directory_path (str): Directory path containing presentation files to merge.
    - output_filename (str): Filename for the output merged presentation.

    Returns:
    - None
    """
    # Create a new presentation to merge the slides into
    merged_presentation = Presentation()
    merged_presentation.slide_width = Inches(PPT_WIDTH)
    merged_presentation.slide_height = Inches(PPT_HEIGHT)
    # Iterate over each presentation file in the directory and copy its slides
    files = sorted(os.listdir(directory_path))
    print(f"Merging {len(files)} slides")
    for filename in files:
        if filename.endswith(".pptx"):
            print(f"adding {filename} to report")
            filepath = os.path.join(directory_path, filename)
            source_presentation = Presentation(filepath)
            for slide in source_presentation.slides:
                copy_slide(slide, merged_presentation)

    output_ppt_path = output_filename.replace('.pptx', '.ppt')
    merged_presentation.save(output_ppt_path)
